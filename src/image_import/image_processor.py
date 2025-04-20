import os
from pathlib import Path
from typing import List, Dict
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt

class ImageProcessor:
    def __init__(self):
        self.image_folder = None
        self.image_files = []
        self.grouped_images = {}
        
    def set_image_folder(self, folder_path: str) -> List[str]:
        """设置图片文件夹并返回所有图片文件"""
        self.image_folder = folder_path
        self.image_files = []
        
        # 支持的图片格式
        image_extensions = ('.jpg', '.jpeg', '.png', '.bmp')
        
        # 遍历文件夹获取所有图片
        for file in os.listdir(folder_path):
            if file.lower().endswith(image_extensions):
                self.image_files.append(os.path.join(folder_path, file))
                
        return self.image_files
    
    def group_images_by_batch(self, batch_pattern: str = None) -> Dict[str, List[str]]:
        """按批次号对图片进行分组"""
        self.grouped_images = {}
        
        for image_path in self.image_files:
            # 从文件名中提取批次号
            filename = os.path.basename(image_path)
            if batch_pattern:
                # 使用正则表达式提取批次号
                import re
                match = re.search(batch_pattern, filename)
                if match:
                    batch_id = match.group(1)
                else:
                    batch_id = "未分组"
            else:
                # 默认使用文件名中的数字作为批次号
                batch_id = ''.join(filter(str.isdigit, filename)) or "未分组"
            
            if batch_id not in self.grouped_images:
                self.grouped_images[batch_id] = []
            self.grouped_images[batch_id].append(image_path)
            
        return self.grouped_images
    
    def create_ppt_report(self, template_path: str = None, output_path: str = "report.pptx"):
        """创建PPT报告"""
        # 创建新的PPT或使用模板
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
            
        # 为每个批次创建一个新的幻灯片
        for batch_id, images in self.grouped_images.items():
            # 添加标题幻灯片
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = f"批次 {batch_id} 检测报告"
            subtitle.text = f"共 {len(images)} 张图片"
            
            # 添加图片幻灯片
            content_slide = prs.slides.add_slide(prs.slide_layouts[5])
            
            # 在幻灯片上添加图片
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            
            for idx, image_path in enumerate(images):
                if idx > 0:  # 每张幻灯片最多放一张图片
                    content_slide = prs.slides.add_slide(prs.slide_layouts[5])
                
                content_slide.shapes.add_picture(
                    image_path,
                    left, top,
                    width=width
                )
                
                # 添加图片说明
                textbox = content_slide.shapes.add_textbox(
                    left, top + width * 0.75,
                    width, Inches(1)
                )
                text_frame = textbox.text_frame
                p = text_frame.add_paragraph()
                p.text = f"图片 {idx + 1}: {os.path.basename(image_path)}"
                p.font.size = Pt(12)
        
        # 保存PPT
        prs.save(output_path)
        return output_path 
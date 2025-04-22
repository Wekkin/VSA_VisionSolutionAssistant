import sys
import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
                             QFileDialog, QListWidget, QPushButton, QLabel, QTextEdit,
                             QGraphicsView, QGraphicsScene, QGraphicsPixmapItem, QGraphicsRectItem,
                             QSplitter)
from PyQt5.QtCore import Qt, QRectF, QPoint
from PyQt5.QtGui import QPixmap, QImage, QPainter, QColor


class ImageProcessor(QGraphicsView):
    def __init__(self):
        super().__init__()
        self.scene = QGraphicsScene()
        self.setScene(self.scene)
        self.crop_rect = None
        self.start_pos = None
        self.current_pixmap = None
        self.cropped_pixmap = None
        self.current_image_path = None
        self.setRenderHint(QPainter.Antialiasing)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.setVerticalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.setViewportUpdateMode(QGraphicsView.FullViewportUpdate)
        self.setRenderHint(QPainter.SmoothPixmapTransform)
        self.crop_completed = None  # 回调函数

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self.adjust_image()
        if self.cropped_pixmap and self.crop_completed:
            self.crop_completed()  # 调用回调函数更新裁剪视图

    def load_image(self, path):
        # 清除之前的裁剪矩形
        if self.crop_rect and self.crop_rect.scene() == self.scene:
            self.scene.removeItem(self.crop_rect)
        self.crop_rect = None
        self.start_pos = None
        
        # 清除场景并加载新图片
        self.scene.clear()
        self.current_image_path = path
        pixmap = QPixmap(str(path))
        self.current_pixmap = pixmap
        self.cropped_pixmap = None
        self.pixmap_item = self.scene.addPixmap(pixmap)
        self.adjust_image()

    def adjust_image(self):
        if not self.current_pixmap:
            return
            
        try:
            # 获取视图和图片的尺寸
            view_rect = self.viewport().rect()
            pixmap_rect = self.current_pixmap.rect()
            
            # 计算缩放比例
            scale_w = view_rect.width() / pixmap_rect.width()
            scale_h = view_rect.height() / pixmap_rect.height()
            scale = min(scale_w, scale_h)  # 使用较小的缩放比例以确保完整显示
            
            # 重置变换
            self.resetTransform()
            # 应用缩放
            self.scale(scale, scale)
            
            # 将QRect转换为QRectF
            scene_rect = QRectF(
                pixmap_rect.x(),
                pixmap_rect.y(),
                pixmap_rect.width(),
                pixmap_rect.height()
            )
            # 调整场景大小
            self.setSceneRect(scene_rect)
            # 居中显示
            self.centerOn(self.pixmap_item)
        except Exception as e:
            print(f"Error in adjust_image: {str(e)}")

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            # 清除之前的裁剪矩形
            if self.crop_rect and self.crop_rect.scene() == self.scene:
                self.scene.removeItem(self.crop_rect)
            
            self.start_pos = self.mapToScene(event.pos())
            self.crop_rect = QGraphicsRectItem()
            self.crop_rect.setPen(QColor(255, 0, 0))
            self.scene.addItem(self.crop_rect)

    def mouseMoveEvent(self, event):
        if self.start_pos and self.crop_rect and self.crop_rect.scene() == self.scene:
            end_pos = self.mapToScene(event.pos())
            rect = QRectF(self.start_pos, end_pos).normalized()
            self.crop_rect.setRect(rect)

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.LeftButton and self.start_pos and self.crop_rect and self.crop_rect.scene() == self.scene:
            end_pos = self.mapToScene(event.pos())
            rect = QRectF(self.start_pos, end_pos).normalized()
            
            # 获取当前缩放比例
            transform = self.transform()
            scale = transform.m11()  # 获取水平缩放因子
            
            # 计算实际图片上的裁剪区域
            actual_rect = QRectF(
                rect.x() / scale,
                rect.y() / scale,
                rect.width() / scale,
                rect.height() / scale
            )
            
            # 确保裁剪区域在图片范围内
            if self.current_pixmap:
                img_rect = self.current_pixmap.rect()
                actual_rect = actual_rect.intersected(QRectF(img_rect))
                
                if not actual_rect.isEmpty():
                    self.crop_area = (
                        actual_rect.left(),
                        actual_rect.top(),
                        actual_rect.right(),
                        actual_rect.bottom()
                    )
                    
                    # 创建裁剪后的图片
                    cropped = self.current_pixmap.copy(
                        int(actual_rect.left()),
                        int(actual_rect.top()),
                        int(actual_rect.width()),
                        int(actual_rect.height())
                    )
                    self.cropped_pixmap = cropped
                    
                    # 更新裁剪视图
                    if self.crop_completed:
                        self.crop_completed()
            
            self.start_pos = None


class DetailViewer(QLabel):
    def __init__(self):
        super().__init__()
        self.setMinimumSize(100, 100)
        self.setAlignment(Qt.AlignCenter)
        self.setStyleSheet("border: 1px solid #cccccc;")

    def resizeEvent(self, event):
        super().resizeEvent(event)
        if self.pixmap():
            self.update_image(self.pixmap())

    def update_image(self, pixmap):
        if pixmap:
            scaled_pixmap = pixmap.scaled(
                self.size(),
                Qt.KeepAspectRatio,
                Qt.SmoothTransformation
            )
            super().setPixmap(scaled_pixmap)


class PPTGeneratorApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.image_data = {}
        self.current_image = None
        self.default_comment = "评估结论：1.成像清晰，检测无风险；2.可以通过控制阀值完成允收"

    def init_ui(self):
        self.setWindowTitle("幻灯片生成器")
        self.setGeometry(100, 100, 1500, 900)

        # 创建主窗口部件
        main_widget = QWidget()
        main_layout = QVBoxLayout(main_widget)

        # 创建内容区域
        content_widget = QWidget()
        content_layout = QHBoxLayout(content_widget)
        content_layout.setContentsMargins(0, 0, 0, 0)

        # 创建主分割器
        main_splitter = QSplitter(Qt.Horizontal)
        main_splitter.setHandleWidth(1)  # 设置分割条的宽度
        main_splitter.setStyleSheet("""
            QSplitter::handle {
                background: #cccccc;
            }
            QSplitter::handle:hover {
                background: #999999;
            }
        """)

        # 左侧文件列表区域
        left_panel = QWidget()
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(5, 5, 5, 5)
        
        file_list_label = QLabel("图片导入")
        file_list_label.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px;")
        
        self.btn_load = QPushButton("选择图片文件夹")
        self.btn_load.clicked.connect(self.load_folders)
        
        self.image_list = QListWidget()
        
        left_layout.addWidget(file_list_label)
        left_layout.addWidget(self.btn_load)
        left_layout.addWidget(self.image_list)

        # 右侧主要内容区域
        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(5, 5, 5, 5)

        # 创建右侧垂直分割器
        right_splitter = QSplitter(Qt.Vertical)
        right_splitter.setHandleWidth(1)

        # PPT标题区域
        title_widget = QWidget()
        title_layout = QHBoxLayout(title_widget)
        title_layout.setContentsMargins(5, 5, 5, 5)
        title_label = QLabel("PPT标题")
        title_label.setStyleSheet("font-size: 15px; font-weight: bold;")
        self.title_edit = QTextEdit()
        self.title_edit.setMaximumHeight(35)
        title_layout.addWidget(title_label)
        title_layout.addWidget(self.title_edit)

        # 图片显示区域的水平分割器
        display_splitter = QSplitter(Qt.Horizontal)
        display_splitter.setHandleWidth(1)
        
        # 原图显示
        original_widget = QWidget()
        original_layout = QVBoxLayout(original_widget)
        original_layout.setContentsMargins(5, 5, 5, 5)
        original_label = QLabel("原图展示")
        original_label.setStyleSheet("font-size: 15px; font-weight: bold; text-align: center;")
        self.image_processor = ImageProcessor()
        original_layout.addWidget(original_label)
        original_layout.addWidget(self.image_processor)
        
        # 细节图显示
        detail_widget = QWidget()
        detail_layout = QVBoxLayout(detail_widget)
        detail_layout.setContentsMargins(5, 5, 5, 5)
        detail_label = QLabel("细节放大")
        detail_label.setStyleSheet("font-size: 15px; font-weight: bold; text-align: center;")
        self.detail_view = DetailViewer()
        detail_layout.addWidget(detail_label)
        detail_layout.addWidget(self.detail_view)

        # 添加到水平分割器
        display_splitter.addWidget(original_widget)
        display_splitter.addWidget(detail_widget)
        display_splitter.setStretchFactor(0, 1)
        display_splitter.setStretchFactor(1, 1)

        # 可行性分析区域
        analysis_widget = QWidget()
        analysis_layout = QVBoxLayout(analysis_widget)
        analysis_layout.setContentsMargins(5, 5, 5, 5)
        analysis_label = QLabel("可行性分析")
        analysis_label.setStyleSheet("font-size: 15px; font-weight: bold;")
        self.comment_edit = QTextEdit()
        self.comment_edit.setPlaceholderText("在此输入评估意见...")
        analysis_layout.addWidget(analysis_label)
        analysis_layout.addWidget(self.comment_edit)

        # 底部按钮区域
        bottom_widget = QWidget()
        bottom_layout = QHBoxLayout(bottom_widget)
        bottom_layout.setContentsMargins(5, 5, 5, 5)
        self.btn_generate = QPushButton("生成PPT")
        self.btn_generate.clicked.connect(self.generate_ppt)
        self.btn_generate.setFixedSize(120, 40)
        self.btn_generate.setStyleSheet("background-color: #4CAF50; color: white; font-weight: bold;")
        bottom_layout.addStretch()
        bottom_layout.addWidget(self.btn_generate)

        # 将组件添加到右侧垂直分割器
        right_splitter.addWidget(title_widget)
        right_splitter.addWidget(display_splitter)
        right_splitter.addWidget(analysis_widget)
        right_splitter.addWidget(bottom_widget)

        # 设置右侧分割器的比例
        right_splitter.setStretchFactor(0, 1)  # 标题
        right_splitter.setStretchFactor(1, 8)  # 图片显示
        right_splitter.setStretchFactor(2, 3)  # 分析
        right_splitter.setStretchFactor(3, 1)  # 底部按钮

        # 将左右面板添加到主分割器
        main_splitter.addWidget(left_panel)
        main_splitter.addWidget(right_splitter)
        
        # 设置左右面板的初始比例
        main_splitter.setStretchFactor(0, 1)  # 左侧面板
        main_splitter.setStretchFactor(1, 4)  # 右侧面板

        # 将主分割器添加到主布局
        content_layout.addWidget(main_splitter)
        main_layout.addWidget(content_widget)
        
        self.setCentralWidget(main_widget)

        # 设置整体样式
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f5f5f5;
            }
            QPushButton {
                padding: 5px;
                border-radius: 3px;
                background-color: #f0f0f0;
            }
            QPushButton:hover {
                background-color: #e0e0e0;
            }
            QListWidget {
                border: 1px solid #cccccc;
                border-radius: 35px;
            }
            QTextEdit {
                border: 1px solid #cccccc;
                border-radius: 3px;
                padding: 5px;
            }
        """)

        # 添加信号连接
        self.image_list.currentItemChanged.connect(self.select_image)
        self.image_processor.crop_completed = self.update_detail_view

    def load_folders(self):
        folders = QFileDialog.getExistingDirectory(self, "选择包含图片的文件夹")
        if folders:
            self.scan_images(Path(folders))

    def scan_images(self, root_path):
        self.image_data.clear()
        self.image_list.clear()

        for img_path in root_path.rglob("*.*"):
            if img_path.suffix.lower() in ['.jpg', '.jpeg', '.png']:
                folder_name = img_path.parent.name
                rel_path = str(img_path.relative_to(root_path))

                self.image_data[rel_path] = {
                    'full_path': str(img_path),
                    'crop_area': None,
                    'comment': self.default_comment,
                    'folder': folder_name
                }
                self.image_list.addItem(rel_path)

    def update_detail_view(self):
        """更新细节视图"""
        if hasattr(self.image_processor, 'cropped_pixmap') and self.image_processor.cropped_pixmap:
            self.detail_view.update_image(self.image_processor.cropped_pixmap)

    def select_image(self, item):
        if item:
            self.current_image = item.text()
            img_info = self.image_data[self.current_image]
            self.image_processor.load_image(img_info['full_path'])
            self.comment_edit.setPlainText(img_info['comment'])
            # 清除细节视图
            self.detail_view.clear()
            # 设置标题
            self.title_edit.setPlainText(img_info['folder'])

    def save_current_data(self):
        if self.current_image:
            img_info = self.image_data[self.current_image]
            # 保存裁剪区域
            if hasattr(self.image_processor, 'crop_area'):
                img_info['crop_area'] = self.image_processor.crop_area
            # 保存评估意见
            img_info['comment'] = self.comment_edit.toPlainText() or self.default_comment
            # 保存标题
            img_info['folder'] = self.title_edit.toPlainText()

    def generate_ppt(self):
        self.save_current_data()
        prs = Presentation()
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
        prs.slide_width = slide_width
        prs.slide_height = slide_height

        for rel_path, data in self.image_data.items():
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 添加标题
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2),
                                           slide_width - Inches(1), Inches(0.8))
            title.text_frame.text = data['folder']

            try:
                # 添加原图（左侧）
                original_img = slide.shapes.add_picture(
                    data['full_path'],
                    Inches(0.5), Inches(1.5),
                    width=Inches(6), height=Inches(4)
                )

                # 添加裁剪图（右侧）
                if data['crop_area']:
                    with Image.open(data['full_path']) as img:
                        cropped = img.crop(data['crop_area'])
                        temp_path = f"temp_{Path(rel_path).name}"
                        cropped.save(temp_path)

                        slide.shapes.add_picture(
                            temp_path,
                            Inches(7), Inches(1.5),
                            width=Inches(5), height=Inches(3.75)
                        )
                        os.remove(temp_path)

                # 添加评估意见
                comment_box = slide.shapes.add_textbox(
                    Inches(0.3), Inches(6),
                    slide_width - Inches(1), Inches(1)
                )
                comment_box.text_frame.text = data['comment']

            except Exception as e:
                print(f"Error processing {rel_path}: {str(e)}")
                continue

        prs.save("AutoGenerated_Presentation.pptx")
        self.statusBar().showMessage("PPT生成完成！", 5000)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    window = PPTGeneratorApp()
    window.show()
    sys.exit(app.exec_())
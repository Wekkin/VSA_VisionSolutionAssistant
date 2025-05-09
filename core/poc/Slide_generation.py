import sys
import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
                             QFileDialog, QListWidget, QPushButton, QLabel, QTextEdit,
                             QGraphicsView, QGraphicsScene, QGraphicsPixmapItem, QGraphicsRectItem,
                             QSplitter, QMessageBox)
from PyQt5.QtCore import Qt, QRectF, QPoint
from PyQt5.QtGui import QPixmap, QImage, QPainter, QColor, QPen
from datetime import datetime
import subprocess
import json
from utils.path_utils import get_resource_path


class ImageProcessor(QGraphicsView):
    def __init__(self):
        super().__init__()
        self.scene = QGraphicsScene()
        self.setScene(self.scene)
        self.crop_rect = None
        self.start_pos = None
        self.current_pixmap = None
        self.cropped_pixmap = None
        self.current_image_path = None
        self.pixmap_item = None
        self.setRenderHint(QPainter.Antialiasing)
        self.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.setVerticalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.setViewportUpdateMode(QGraphicsView.FullViewportUpdate)
        self.setRenderHint(QPainter.SmoothPixmapTransform)
        self.crop_completed = None
        self.crop_area = None
        self.scale_factor = 1.0
        
        # 设置场景背景
        self.setStyleSheet("""
            QGraphicsView {
                border: 1px solid #cccccc;
                background: white;
            }
        """)

    def resizeEvent(self, event):
        """窗口大小改变事件"""
        super().resizeEvent(event)
        self.adjust_image()
        if self.cropped_pixmap and self.crop_completed:
            self.crop_completed()  # 调用回调函数更新裁剪视图

    def load_image(self, path):
        """加载图片"""
        try:
            # 清除现有内容
            self.scene.clear()
            self.crop_rect = None
            self.start_pos = None
            self.current_image_path = path
            self.cropped_pixmap = None
            self.crop_area = None
            
            # 加载图片
            self.current_pixmap = QPixmap(str(path))
            if not self.current_pixmap.isNull():
                self.pixmap_item = self.scene.addPixmap(self.current_pixmap)
                self.adjust_image()
                print(f"Image loaded successfully: {path}")
            else:
                print(f"Failed to load image: {path}")
        except Exception as e:
            print(f"Error loading image: {str(e)}")

    def adjust_image(self):
        """调整图片大小和位置"""
        if not self.current_pixmap or self.current_pixmap.isNull():
            return
            
        try:
            # 获取视图和图片的尺寸
            view_rect = self.viewport().rect()
            pixmap_rect = self.current_pixmap.rect()
            
            # 计算缩放比例
            scale_w = view_rect.width() / pixmap_rect.width()
            scale_h = view_rect.height() / pixmap_rect.height()
            self.scale_factor = min(scale_w, scale_h)
            
            # 重置变换
            self.resetTransform()
            # 应用缩放
            self.scale(self.scale_factor, self.scale_factor)
            
            # 设置场景大小
            self.scene.setSceneRect(0, 0, pixmap_rect.width(), pixmap_rect.height())
            # 居中显示
            self.centerOn(self.pixmap_item)
            
            print(f"Image adjusted, scale: {self.scale_factor}")
        except Exception as e:
            print(f"Error adjusting image: {str(e)}")

    def mousePressEvent(self, event):
        """鼠标按下事件"""
        if event.button() == Qt.LeftButton and self.pixmap_item:
            try:
                # 获取场景坐标
                scene_pos = self.mapToScene(event.pos())
                # 转换为图片坐标
                self.start_pos = scene_pos
                
                # 移除现有的选择框
                if self.crop_rect:
                    self.scene.removeItem(self.crop_rect)
                
                # 创建新的选择框
                self.crop_rect = QGraphicsRectItem()
                pen = QPen(QColor(255, 0, 0))  # 红色
                pen.setWidth(2)  # 设置线宽
                self.crop_rect.setPen(pen)
                self.scene.addItem(self.crop_rect)
                
                print(f"Mouse press at: {scene_pos}")
            except Exception as e:
                print(f"Error in mouse press: {str(e)}")

    def mouseMoveEvent(self, event):
        """鼠标移动事件"""
        if event.buttons() & Qt.LeftButton and self.start_pos and self.crop_rect:
            try:
                # 获取当前位置
                current_pos = self.mapToScene(event.pos())
                
                # 创建选择区域
                rect = QRectF(self.start_pos, current_pos).normalized()
                self.crop_rect.setRect(rect)
                
                print(f"Mouse move to: {current_pos}")
            except Exception as e:
                print(f"Error in mouse move: {str(e)}")

    def mouseReleaseEvent(self, event):
        """鼠标释放事件"""
        if event.button() == Qt.LeftButton and self.start_pos and self.crop_rect:
            try:
                # 获取结束位置
                end_pos = self.mapToScene(event.pos())
                rect = QRectF(self.start_pos, end_pos).normalized()
                
                # 确保裁剪区域在图片范围内
                if self.current_pixmap:
                    img_rect = QRectF(self.current_pixmap.rect())
                    actual_rect = rect.intersected(img_rect)
                    
                    if not actual_rect.isEmpty():
                        # 保存裁剪区域
                        self.crop_area = (
                            int(actual_rect.x()),
                            int(actual_rect.y()),
                            int(actual_rect.width()),
                            int(actual_rect.height())
                        )
                        
                        # 创建裁剪后的图片
                        cropped = self.current_pixmap.copy(
                            int(actual_rect.x()),
                            int(actual_rect.y()),
                            int(actual_rect.width()),
                            int(actual_rect.height())
                        )
                        self.cropped_pixmap = cropped
                        
                        print(f"Crop completed: area={self.crop_area}, size={cropped.size()}")
                        
                        # 调用回调函数更新预览
                        if self.crop_completed:
                            self.crop_completed()
                
                self.start_pos = None
            except Exception as e:
                print(f"Error in mouse release: {str(e)}")


class DetailViewer(QLabel):
    def __init__(self):
        super().__init__()
        self.setMinimumSize(200, 200)
        self.setAlignment(Qt.AlignCenter)
        self.setStyleSheet("""
            QLabel {
                border: 1px solid #cccccc;
                background: white;
            }
        """)
        self._current_pixmap = None

    def clear(self):
        """清除图片"""
        self._current_pixmap = None
        super().clear()
        print("Detail view cleared")

    def update_image(self, pixmap):
        """更新预览图片"""
        try:
            if pixmap and not pixmap.isNull():
                self._current_pixmap = pixmap
                # 计算缩放后的大小
                available_size = self.size()
                scaled_pixmap = pixmap.scaled(
                    available_size,
                    Qt.KeepAspectRatio,
                    Qt.SmoothTransformation
                )
                super().setPixmap(scaled_pixmap)
                print(f"Detail view updated with image size: {pixmap.size()}")
            else:
                print("Warning: Received null or invalid pixmap")
                self.clear()
        except Exception as e:
            print(f"Error updating detail view: {str(e)}")

    def resizeEvent(self, event):
        """处理窗口大小改变事件"""
        super().resizeEvent(event)
        if self._current_pixmap:
            self.update_image(self._current_pixmap)


class PPTGeneratorApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.init_ui()
        self.image_data = {}
        self.current_image = None
        self.default_comment = "评估结论：1.成像清晰，检测无风险；2.可以通过控制阀值完成允收"
        self.cache_dir = None  # 缓存目录
        
        # 设置图片处理器的回调
        self.image_processor.crop_completed = self.update_detail_view

    def init_ui(self):
        self.setWindowTitle("幻灯片生成器")
        self.setGeometry(100, 100, 1500, 900)

        # 创建主窗口部件
        main_widget = QWidget()
        main_layout = QVBoxLayout(main_widget)

        # 创建内容区域
        content_widget = QWidget()
        content_layout = QHBoxLayout(content_widget)
        content_layout.setContentsMargins(0, 0, 0, 0)

        # 创建主分割器
        main_splitter = QSplitter(Qt.Horizontal)
        main_splitter.setHandleWidth(1)  # 设置分割条的宽度
        main_splitter.setStyleSheet("""
            QSplitter::handle {
                background: #cccccc;
            }
            QSplitter::handle:hover {
                background: #999999;
            }
        """)

        # 左侧文件列表区域
        left_panel = QWidget()
        left_layout = QVBoxLayout(left_panel)
        left_layout.setContentsMargins(5, 5, 5, 5)
        
        file_list_label = QLabel("图片导入")
        file_list_label.setStyleSheet("font-size: 15px; font-weight: bold; padding: 5px;")
        
        self.btn_load = QPushButton("选择图片文件夹")
        self.btn_load.clicked.connect(self.load_folders)
        self.btn_load.setStyleSheet("""
            QPushButton {
                background: #1890ff;
                color: white;
                padding: 4px 16px;
                border-radius: 4px;
                border: none;
            }
            QPushButton:hover {
                background: #40a9ff;
            }
        """)
        
        self.image_list = QListWidget()
        
        left_layout.addWidget(file_list_label)
        left_layout.addWidget(self.btn_load)
        left_layout.addWidget(self.image_list)

        # 右侧主要内容区域
        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(5, 5, 5, 5)

        # 创建右侧垂直分割器
        right_splitter = QSplitter(Qt.Vertical)
        right_splitter.setHandleWidth(1)

        # PPT标题区域
        title_widget = QWidget()
        title_layout = QHBoxLayout(title_widget)
        title_layout.setContentsMargins(5, 5, 5, 5)
        title_label = QLabel("PPT标题")
        title_label.setStyleSheet("font-size: 15px; font-weight: bold;")
        self.title_edit = QTextEdit()
        self.title_edit.setMaximumHeight(35)
        self.title_edit.textChanged.connect(self.on_title_changed)  # 添加文本变化监听
        
        # 创建生成PPT按钮
        self.btn_generate = QPushButton("生成PPT")
        self.btn_generate.clicked.connect(self.generate_ppt)
        self.btn_generate.setFixedSize(120, 35)  # 调整按钮高度与标题框一致
        self.btn_generate.setStyleSheet("""
            QPushButton {
                background-color: #4CAF50;
                color: white;
                font-weight: bold;
                border-radius: 4px;
                padding: 5px;
            }
            QPushButton:hover {
                background-color: #45a049;
            }
        """)
        
        # 将组件添加到标题布局中
        title_layout.addWidget(title_label)
        title_layout.addWidget(self.title_edit, stretch=1)  # 让标题框占据更多空间
        title_layout.addWidget(self.btn_generate)

        # 图片显示区域的水平分割器
        display_splitter = QSplitter(Qt.Horizontal)
        display_splitter.setHandleWidth(1)
        
        # 原图显示
        original_widget = QWidget()
        original_layout = QVBoxLayout(original_widget)
        original_layout.setContentsMargins(5, 5, 5, 5)
        original_label = QLabel("原图展示")
        original_label.setStyleSheet("font-size: 15px; font-weight: bold; text-align: center;")
        self.image_processor = ImageProcessor()
        original_layout.addWidget(original_label)
        original_layout.addWidget(self.image_processor)
        
        # 细节图显示
        detail_widget = QWidget()
        detail_layout = QVBoxLayout(detail_widget)
        detail_layout.setContentsMargins(5, 5, 5, 5)
        detail_label = QLabel("细节放大")
        detail_label.setStyleSheet("font-size: 15px; font-weight: bold; text-align: center;")
        self.detail_view = DetailViewer()
        detail_layout.addWidget(detail_label)
        detail_layout.addWidget(self.detail_view)

        # 添加到水平分割器
        display_splitter.addWidget(original_widget)
        display_splitter.addWidget(detail_widget)
        display_splitter.setStretchFactor(0, 1)  # 设置原图展示框的拉伸因子为1
        display_splitter.setStretchFactor(1, 1)  # 设置细节放大框的拉伸因子为1
        display_splitter.setSizes([500, 500])  # 设置初始大小相等

        # 可行性分析区域
        analysis_widget = QWidget()
        analysis_layout = QVBoxLayout(analysis_widget)
        analysis_layout.setContentsMargins(5, 5, 5, 5)
        analysis_label = QLabel("可行性分析")
        analysis_label.setStyleSheet("font-size: 15px; font-weight: bold;")
        self.comment_edit = QTextEdit()
        self.comment_edit.setPlaceholderText("在此输入评估意见...")
        self.comment_edit.textChanged.connect(self.on_comment_changed)  # 添加文本变化监听
        analysis_layout.addWidget(analysis_label)
        analysis_layout.addWidget(self.comment_edit)

        # 将组件添加到右侧垂直分割器
        right_splitter.addWidget(title_widget)
        right_splitter.addWidget(display_splitter)
        right_splitter.addWidget(analysis_widget)

        # 设置右侧分割器的比例
        right_splitter.setStretchFactor(0, 1)  # 标题
        right_splitter.setStretchFactor(1, 8)  # 图片显示
        right_splitter.setStretchFactor(2, 3)  # 分析

        # 将左右面板添加到主分割器
        main_splitter.addWidget(left_panel)
        main_splitter.addWidget(right_splitter)
        
        # 设置左右面板的初始比例
        main_splitter.setStretchFactor(0, 1)  # 左侧面板
        main_splitter.setStretchFactor(1, 4)  # 右侧面板

        # 将主分割器添加到主布局
        content_layout.addWidget(main_splitter)
        main_layout.addWidget(content_widget)
        
        self.setCentralWidget(main_widget)

        # 设置整体样式
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f5f5f5;
            }
            QPushButton {
                padding: 5px;
                border-radius: 3px;
                background-color: #f0f0f0;
            }
            QPushButton:hover {
                background-color: #e0e0e0;
            }
            QListWidget {
                border: 1px solid #cccccc;
                border-radius: 35px;
            }
            QTextEdit {
                border: 1px solid #cccccc;
                border-radius: 3px;
                padding: 5px;
            }
        """)

        # 添加信号连接
        self.image_list.currentItemChanged.connect(self.select_image)

    def load_folders(self):
        folders = QFileDialog.getExistingDirectory(self, "选择包含图片的文件夹")
        if folders:
            self.project_root = Path(folders)
            # 创建缓存目录
            self.cache_dir = self.project_root / '.vsa_cache'
            self.cache_dir.mkdir(exist_ok=True)
            # 创建或加载JSON记录文件
            self.record_file = self.cache_dir / 'crop_records.json'
            self.load_records()
            self.scan_images(self.project_root)

    def load_records(self):
        """加载已有的裁剪记录"""
        try:
            if self.record_file.exists():
                with open(self.record_file, 'r', encoding='utf-8') as f:
                    records = json.load(f)
                    # 更新现有记录
                    for img_name, data in records.items():
                        if img_name in self.image_data:
                            self.image_data[img_name].update(data)
        except Exception as e:
            print(f"Error loading records: {str(e)}")

    def save_records(self):
        """保存记录到JSON文件"""
        try:
            if self.cache_dir and self.record_file:
                records = {}
                for img_name, data in self.image_data.items():
                    records[img_name] = {
                        'crop_area': data.get('crop_area'),
                        'comment': data.get('comment', self.default_comment),
                        'folder': data.get('folder', ''),
                        'folder_name': data.get('folder_name', ''),
                        'cache_path': data.get('cache_path')
                    }
                with open(self.record_file, 'w', encoding='utf-8') as f:
                    json.dump(records, f, indent=2, ensure_ascii=False)
                print("Records saved successfully")
        except Exception as e:
            print(f"Error saving records: {str(e)}")

    def scan_images(self, root_path):
        """扫描项目路径下的所有图片"""
        try:
            # 清除现有数据
            self.image_data.clear()
            self.image_list.clear()
            self.image_processor.scene.clear()
            self.detail_view.clear()
            self.comment_edit.clear()

            # 支持的图片格式
            image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.gif')
            
            # 递归扫描所有图片
            for file_path in Path(root_path).rglob('*'):
                if file_path.suffix.lower() in image_extensions:
                    file_name = file_path.name
                    folder_name = file_path.parent.name
                    
                    # 存储图片信息，使用绝对路径
                    abs_path = str(file_path.absolute())
                    self.image_data[file_name] = {
                        'full_path': abs_path,
                        'folder': folder_name,  # 默认使用文件夹名称
                        'folder_name': folder_name,  # 保存原始文件夹名称
                        'comment': self.default_comment,
                        'crop_area': None
                    }
                    
                    # 添加到列表
                    self.image_list.addItem(file_name)
            
            # 如果找到图片，选择第一张
            if self.image_list.count() > 0:
                self.image_list.setCurrentRow(0)
                
            return True
        except Exception as e:
            print(f"Error scanning images: {str(e)}")
            return False

    def update_detail_view(self):
        """更新细节视图并缓存裁剪图片"""
        try:
            if hasattr(self.image_processor, 'cropped_pixmap') and self.image_processor.cropped_pixmap:
                cropped = self.image_processor.cropped_pixmap
                if not cropped.isNull():
                    # 更新预览
                    self.detail_view.update_image(cropped)
                    
                    # 保存裁剪图到缓存
                    if self.current_image and self.cache_dir:
                        cache_name = f"crop_{self.current_image}"
                        cache_path = self.cache_dir / cache_name
                        cropped.save(str(cache_path))
                        
                        # 更新记录
                        if self.current_image in self.image_data:
                            self.image_data[self.current_image]['cache_path'] = str(cache_path)
                            self.image_data[self.current_image]['crop_area'] = self.image_processor.crop_area
                            # 保存记录到JSON
                            self.save_records()
                            
                    print(f"Updated and cached detail view: {cropped.size()}")
                else:
                    print("Warning: Cropped pixmap is null")
            else:
                print("Warning: No cropped image available")
        except Exception as e:
            print(f"Error updating detail view: {str(e)}")

    def select_image(self, item):
        """选择图片时的处理"""
        if item:
            file_name = item.text()
            if file_name in self.image_data:
                try:
                    self.current_image = file_name
                    img_info = self.image_data[file_name]
                    
                    # 加载图片
                    self.image_processor.load_image(img_info['full_path'])
                    
                    # 设置评估意见
                    self.comment_edit.setPlainText(img_info.get('comment', self.default_comment))
                    
                    # 设置标题（使用文件夹名称作为默认值）
                    folder_name = img_info.get('folder', '')
                    if not folder_name:  # 如果没有保存的标题，使用文件夹名称
                        folder_name = img_info.get('folder_name', '')
                    self.title_edit.setPlainText(folder_name)
                    
                    # 清除细节视图
                    self.detail_view.clear()
                    
                    print(f"Loading image: {img_info['full_path']}")
                except Exception as e:
                    print(f"Error selecting image: {str(e)}")

    def on_title_changed(self):
        """标题文本变化时的处理"""
        if self.current_image:
            self.image_data[self.current_image]['folder'] = self.title_edit.toPlainText()
            self.save_records()

    def on_comment_changed(self):
        """评估意见文本变化时的处理"""
        if self.current_image:
            self.image_data[self.current_image]['comment'] = self.comment_edit.toPlainText()
            self.save_records()

    def generate_ppt(self):
        """生成PPT，使用缓存的裁剪图片"""
        # 使用模板文件创建演示文稿
        template_path = get_resource_path('src/template.pptx')
        if not os.path.exists(template_path):
            QMessageBox.warning(self, "警告", "未找到模板文件！")
            return
            
        # 检查是否已选择项目目录
        if not hasattr(self, 'project_root'):
            QMessageBox.warning(self, "警告", "请先选择项目文件夹！")
            return

        try:
            prs = Presentation(template_path)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 处理每张图片
            for img_name, data in self.image_data.items():
                # 创建新幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[10])

                # 添加标题
                if slide.shapes.title is None:
                    title = slide.shapes.add_textbox(
                        Inches(2.3), Inches(0.5),
                        slide_width - Inches(3), Inches(0.5)
                    )
                    title.text_frame.text = data['folder']
                    # 设置标题字体大小和加粗
                    title.text_frame.paragraphs[0].font.size = Pt(28)  # 增加字体大小
                    title.text_frame.paragraphs[0].font.bold = False    # 设置为粗体
                else:
                    slide.shapes.title.text = data['folder']
                    # 设置标题字体大小和加粗
                    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)  # 增加字体大小
                    slide.shapes.title.text_frame.paragraphs[0].font.bold = True    # 设置为粗体

                try:
                    # 添加原图
                    with Image.open(data['full_path']) as img:
                        img_width, img_height = img.size
                        target_height = Inches(3)
                        target_width = target_height * (img_width / img_height)
                        left = Inches(0.5)
                        top = Inches(1.8)
                        
                        slide.shapes.add_picture(
                            data['full_path'],
                            left, top,
                            height=target_height
                        )

                    # 添加裁剪图
                    cache_path = data.get('cache_path')
                    if cache_path and os.path.exists(cache_path):
                        with Image.open(cache_path) as cropped:
                            crop_width, crop_height = cropped.size
                            crop_target_height = Inches(3)  # 调整为3英寸，与原图一致
                            crop_target_width = crop_target_height * (crop_width / crop_height)
                            crop_left = slide_width - crop_target_width - Inches(1.5)
                            crop_top = top
                            
                            slide.shapes.add_picture(
                                cache_path,
                                crop_left, crop_top,
                                height=crop_target_height
                            )

                    # 添加评估意见
                    comment_added = False
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape.text_frame.text == "评估意见":
                            shape.text_frame.text = data['comment']
                            comment_added = True
                            break
                    
                    if not comment_added:
                        comment_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(5.2),
                            slide_width - Inches(1), Inches(0.3)
                        )
                        comment_box.text_frame.text = data['comment']

                except Exception as e:
                    print(f"Error processing {img_name}: {str(e)}")
                    continue

            # 生成输出文件名
            project_name = self.project_root.name
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"{project_name}_Report_{timestamp}.pptx"
            output_path = self.project_root / output_filename

            # 保存PPT
            prs.save(str(output_path))
            self.statusBar().showMessage(f"PPT生成完成！保存至：{output_path}", 5000)
            
            # 询问是否打开文件
            reply = QMessageBox.question(self, '完成', 
                                       f'PPT已生成到：\n{output_path}\n\n是否立即打开？',
                                       QMessageBox.Yes | QMessageBox.No)
            
            if reply == QMessageBox.Yes:
                if sys.platform == 'win32':
                    os.startfile(str(output_path))
                elif sys.platform == 'darwin':
                    subprocess.run(['open', str(output_path)])
                else:
                    subprocess.run(['xdg-open', str(output_path)])

        except Exception as e:
            QMessageBox.critical(self, "错误", f"生成PPT时发生错误：\n{str(e)}")
            print(f"Error generating PPT: {str(e)}")

    def closeEvent(self, event):
        """程序关闭时清理缓存"""
        try:
            if hasattr(self, 'cache_dir') and self.cache_dir.exists():
                import shutil
                shutil.rmtree(self.cache_dir)
        except Exception as e:
            print(f"Error cleaning cache: {str(e)}")
        event.accept()


if __name__ == '__main__':
    app = QApplication(sys.argv)
    window = PPTGeneratorApp()
    window.show()
    sys.exit(app.exec_())